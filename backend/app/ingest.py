"""Week 1 — Document ingestion pipeline.

Walks a folder of PDFs (and .docx), extracts text per page, strips
repeating header/footer lines, normalizes whitespace, and emits
overlapping chunks as JSONL.

Usage:
    python -m app.ingest                    # uses defaults from config.py
    python -m app.ingest <raw_dir> <out>    # override paths
"""
from __future__ import annotations

import argparse
import json
import re
import sys
import uuid
from collections import Counter
from pathlib import Path
from typing import Iterator

from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.config import (
    RAW_DIR,
    CHUNKS_PATH,
    CHUNK_SIZE,
    CHUNK_OVERLAP,
    MIN_PAGE_CHARS,
    HEADER_FOOTER_THRESHOLD,
)


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_pdf_pages(path: Path) -> list[tuple[int, str]]:
    """Return [(page_number, raw_text), ...] for a PDF. Page numbers are 1-indexed."""
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as e:
            print(f"  ! page {i} extract failed: {e}", file=sys.stderr)
            text = ""
        pages.append((i, text))
    return pages


def extract_docx_pages(path: Path) -> list[tuple[int, str]]:
    """DOCX has no real pages — return the whole doc as one 'page' (page=1)."""
    from docx import Document  # imported lazily so PDF-only runs don't need it
    doc = Document(str(path))
    text = "\n".join(p.text for p in doc.paragraphs)
    return [(1, text)]


def extract_pptx_pages(path: Path) -> list[tuple[int, str]]:
    """Each slide is one 'page'. Walks all shapes on each slide and collects text."""
    from pptx import Presentation  # lazy import
    prs = Presentation(str(path))
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            # Most text lives in shape.text_frame; tables and other shapes need different handling
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        # Speaker notes — often contain the actual lecture content beyond bullet points
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(notes)
        pages.append((i, "\n".join(parts)))
    return pages


# ---------------------------------------------------------------------------
# Cleaning
# ---------------------------------------------------------------------------

# Strip trailing page counters like "2 / 26", "21 of 30", "page 5" so that
# otherwise-identical header/footer lines collapse to the same key for dedup.
_TRAILING_COUNTER = re.compile(
    r"\s*(page\s+)?\d+\s*(/|of)\s*\d+\s*$|\s*page\s+\d+\s*$",
    re.IGNORECASE,
)


def _normalize_for_dedup(line: str) -> str:
    return _TRAILING_COUNTER.sub("", line).strip()


def find_repeating_lines(pages: list[tuple[int, str]], threshold: float) -> set[str]:
    """Identify lines that appear on >threshold fraction of pages — likely headers/footers."""
    if len(pages) < 3:
        return set()

    line_counts: Counter[str] = Counter()
    for _, text in pages:
        seen_on_page = set()
        for ln in text.splitlines():
            norm = _normalize_for_dedup(ln)
            if len(norm) > 3:
                seen_on_page.add(norm)
        line_counts.update(seen_on_page)

    cutoff = max(2, int(len(pages) * threshold))
    return {line for line, count in line_counts.items() if count >= cutoff}


def normalize_text(text: str, drop_lines: set[str]) -> str:
    """Drop repeating header/footer lines, collapse whitespace."""
    kept = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        # Compare against drop_lines using the same normalization used to build them
        if _normalize_for_dedup(stripped) in drop_lines:
            continue
        # Drop lines that are only page numbers (e.g., "12", "Page 12 of 40")
        if re.fullmatch(r"(page\s+)?\d+(\s*/\s*\d+|\s+of\s+\d+)?", stripped, re.IGNORECASE):
            continue
        kept.append(stripped)

    joined = "\n".join(kept)
    # Collapse 3+ newlines to 2 (paragraph break)
    joined = re.sub(r"\n{3,}", "\n\n", joined)
    # Collapse runs of spaces
    joined = re.sub(r"[ \t]+", " ", joined)
    return joined.strip()


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def make_splitter() -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )


def chunk_document(source_name: str, pages: list[tuple[int, str]]) -> Iterator[dict]:
    """Yield chunk dicts for one document."""
    drop = find_repeating_lines(pages, HEADER_FOOTER_THRESHOLD)
    splitter = make_splitter()

    for page_num, raw in pages:
        cleaned = normalize_text(raw, drop)
        if len(cleaned) < MIN_PAGE_CHARS:
            continue
        for piece in splitter.split_text(cleaned):
            piece = piece.strip()
            if len(piece) < MIN_PAGE_CHARS:
                continue
            yield {
                "id": str(uuid.uuid4()),
                "source": source_name,
                "page": page_num,
                "text": piece,
                "char_count": len(piece),
            }


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------

SUPPORTED_EXT = {".pdf", ".docx", ".pptx"}


def iter_source_files(raw_dir: Path) -> Iterator[Path]:
    for path in sorted(raw_dir.rglob("*")):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXT:
            yield path


def run(raw_dir: Path, out_path: Path) -> None:
    if not raw_dir.exists():
        print(f"ERROR: raw_dir does not exist: {raw_dir}", file=sys.stderr)
        sys.exit(1)

    out_path.parent.mkdir(parents=True, exist_ok=True)

    file_count = 0
    chunk_count = 0
    with out_path.open("w", encoding="utf-8") as f:
        for src in iter_source_files(raw_dir):
            print(f"-> {src.name}")
            try:
                ext = src.suffix.lower()
                if ext == ".pdf":
                    pages = extract_pdf_pages(src)
                elif ext == ".pptx":
                    pages = extract_pptx_pages(src)
                else:
                    pages = extract_docx_pages(src)
            except Exception as e:
                print(f"  ! failed to read {src.name}: {e}", file=sys.stderr)
                continue

            file_chunks = 0
            for chunk in chunk_document(src.name, pages):
                f.write(json.dumps(chunk, ensure_ascii=False) + "\n")
                file_chunks += 1

            print(f"   {len(pages)} pages -> {file_chunks} chunks")
            file_count += 1
            chunk_count += file_chunks

    if file_count == 0:
        print(f"\nNo .pdf or .docx files found in {raw_dir}", file=sys.stderr)
        sys.exit(1)

    print(f"\nDone. Processed {file_count} files, wrote {chunk_count} chunks to {out_path}")


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Ingest PDFs/DOCX into chunked JSONL.")
    p.add_argument("raw_dir", nargs="?", default=str(RAW_DIR), help="Folder of source docs")
    p.add_argument("out_path", nargs="?", default=str(CHUNKS_PATH), help="Output JSONL path")
    return p.parse_args()


if __name__ == "__main__":
    args = parse_args()
    run(Path(args.raw_dir), Path(args.out_path))
